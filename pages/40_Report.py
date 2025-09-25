"""Report export page for PDF / Excel / Word outputs."""
from __future__ import annotations

import io
from dataclasses import dataclass
from datetime import datetime
from decimal import Decimal
from typing import Callable, Dict, List, Sequence, Tuple, TypeVar

import pandas as pd
import streamlit as st
from docx import Document
from docx.shared import Inches
import matplotlib.pyplot as plt
from openpyxl.drawing.image import Image as ExcelImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches as PptxInches, Pt as PptxPt
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import (
    Image as PlatypusImage,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from ui.streamlit_compat import use_container_width_kwargs

from calc import compute, generate_cash_flow, plan_from_models, summarize_plan_metrics
from formatting import format_amount_with_unit, format_ratio
from state import ensure_session_defaults, load_finance_bundle
from theme import THEME_COLORS, inject_theme
from ui.navigation import render_global_navigation, render_workflow_banner

st.set_page_config(
    page_title="経営計画スタジオ｜レポート",
    page_icon="報",
    layout="wide",
)

inject_theme()
ensure_session_defaults()

render_global_navigation("report")
render_workflow_banner("report")

settings_state: Dict[str, object] = st.session_state.get("finance_settings", {})
unit = str(settings_state.get("unit", "百万円"))
fte = Decimal(str(settings_state.get("fte", 20)))
fiscal_year = int(settings_state.get("fiscal_year", 2025))

bundle, has_custom_inputs = load_finance_bundle()
if not has_custom_inputs:
    st.info("入力データが未保存のため、既定値でレポートを生成します。")

plan_cfg = plan_from_models(
    bundle.sales,
    bundle.costs,
    bundle.capex,
    bundle.loans,
    bundle.tax,
    fte=fte,
    unit=unit,
)

amounts = compute(plan_cfg)
metrics = summarize_plan_metrics(amounts)
cash_flow_data = generate_cash_flow(amounts, bundle.capex, bundle.loans, bundle.tax)

st.title("レポート出力")
st.caption("McKinseyスタイルのテンプレートで、主要指標とKPIをPDF / PowerPoint / Excel / Word形式に出力できます。")

SUPPORT_CONTACT = "support@keieiplan.jp"
PDF_FONT_NAME = "HeiseiKakuGo-W5"
T = TypeVar("T")


@dataclass(frozen=True)
class ReportSectionMeta:
    key: str
    title: str
    description: str
    note: str | None = None
    aria_label: str | None = None


REPORT_SECTION_ORDER: Tuple[ReportSectionMeta, ...] = (
    ReportSectionMeta(
        key="エグゼクティブサマリー",
        title="エグゼクティブサマリー",
        description="主要KPIと意思決定ポイントを俯瞰し、経営陣向けの要約として活用できます。",
        note="FTEと表示単位はレポートサブタイトル記載の値を参照してください。",
        aria_label="エグゼクティブサマリー表",
    ),
    ReportSectionMeta(
        key="現状分析",
        title="現状分析（損益計画サマリー）",
        description="損益計画を政策金融機関フォーマットに合わせて整理しています。",
        note="売上高から税引後利益までを標準項目と注記で確認できます。",
        aria_label="現状分析テーブル",
    ),
    ReportSectionMeta(
        key="施策提案",
        title="施策提案（推奨アクションプラン）",
        description="AIがKPIから導出した推奨施策と優先度を提示します。",
        note="優先度は粗利率・キャッシュフロー・損益分岐点の余裕をもとにスコアリングしています。",
        aria_label="施策提案一覧",
    ),
    ReportSectionMeta(
        key="財務シミュレーション",
        title="財務シミュレーション（キャッシュ・投資計画）",
        description="キャッシュフロー区分と投資回収の指標をまとめ、金融機関との対話資料に利用できます。",
        note="投資回収期間およびNPVは初年度キャッシュフローと割引率を基に算出しています。",
        aria_label="財務シミュレーション表",
    ),
)
REPORT_SECTION_LOOKUP: Dict[str, ReportSectionMeta] = {
    section.key: section for section in REPORT_SECTION_ORDER
}
REPORT_SECTIONS: Sequence[str] = tuple(section.key for section in REPORT_SECTION_ORDER)


def _normalize_decimal(value: Decimal | None) -> Decimal:
    if value is None:
        return Decimal("0")
    if isinstance(value, Decimal):
        return Decimal("0") if value.is_nan() else value
    try:
        return Decimal(str(value))
    except Exception:
        return Decimal("0")


def _safe_ratio(value: Decimal, base: Decimal) -> str:
    if base and base != 0:
        return format_ratio(value / base)
    return "—"


def _build_executive_summary_table(
    amounts_data: Dict[str, Decimal],
    metrics_data: Dict[str, Decimal],
    cash_data: Dict[str, Decimal],
    unit: str,
    fte: Decimal,
) -> List[List[str]]:
    sales = _normalize_decimal(amounts_data.get("REV"))
    gross_margin = _normalize_decimal(metrics_data.get("gross_margin"))
    operating_profit = _normalize_decimal(amounts_data.get("OP"))
    operating_margin = _normalize_decimal(metrics_data.get("op_margin"))
    operating_cf = _normalize_decimal(cash_data.get("営業キャッシュフロー"))
    investment_cf = _normalize_decimal(cash_data.get("投資キャッシュフロー"))
    breakeven = _normalize_decimal(metrics_data.get("breakeven"))
    revenue_per_fte_text = "FTE設定が必要です。"
    if fte and fte > 0 and sales:
        per_fte = sales / fte
        revenue_per_fte_text = f"FTEあたり {format_amount_with_unit(per_fte, unit)}"
    margin_highlight = (
        "粗利率45%以上で投資余力を確保。"
        if gross_margin >= Decimal("0.45")
        else "粗利率を底上げするために原価・値付けの再設計を検討。"
        if gross_margin < Decimal("0.35")
        else "チャネルミックス再編で粗利率の更なる改善余地あり。"
    )
    op_highlight = (
        "営業利益率12%以上で収益性は良好。"
        if operating_margin >= Decimal("0.12")
        else "販管費効率化で営業利益率の改善余地あり。"
        if operating_margin < Decimal("0.08")
        else "営業利益率は安定レンジ。追加投資の吸収余力を検証。"
    )
    net_after_invest = operating_cf + investment_cf
    cash_highlight = (
        "営業CFで投資キャッシュを十分にカバーできています。"
        if net_after_invest >= Decimal("0")
        else "営業CFだけでは投資を賄えないため調達計画の検討が必要です。"
    )
    breakeven_display = "—"
    if breakeven > 0:
        breakeven_display = format_amount_with_unit(breakeven, unit)
    breakeven_gap = sales - breakeven if breakeven > 0 else Decimal("0")
    if breakeven > 0:
        breakeven_highlight = (
            f"損益分岐点比 {format_amount_with_unit(breakeven_gap, unit)}の安全余裕。"
            if breakeven_gap >= 0
            else f"損益分岐点まで {format_amount_with_unit(abs(breakeven_gap), unit)}の売上強化が必要。"
        )
    else:
        breakeven_highlight = "損益分岐点の算出には追加データが必要です。"
    return [
        ["指標", "値", "ハイライト"],
        ["年間売上高", format_amount_with_unit(sales, unit), revenue_per_fte_text],
        ["粗利率", format_ratio(gross_margin), margin_highlight],
        ["営業利益", format_amount_with_unit(operating_profit, unit), op_highlight],
        ["営業キャッシュフロー", format_amount_with_unit(operating_cf, unit), cash_highlight],
        ["損益分岐点売上高", breakeven_display, breakeven_highlight],
    ]


def _build_current_analysis_table(
    amounts_data: Dict[str, Decimal],
    metrics_data: Dict[str, Decimal],
    cash_data: Dict[str, Decimal],
    unit: str,
) -> List[List[str]]:
    sales = _normalize_decimal(amounts_data.get("REV"))
    cogs = _normalize_decimal(amounts_data.get("COGS_TTL"))
    gross = _normalize_decimal(amounts_data.get("GROSS"))
    opex = _normalize_decimal(amounts_data.get("OPEX_TTL"))
    op = _normalize_decimal(amounts_data.get("OP"))
    ord_income = _normalize_decimal(amounts_data.get("ORD"))
    non_operating_income = sum(
        _normalize_decimal(amounts_data.get(code)) for code in ("NOI_MISC", "NOI_GRANT", "NOI_OTH")
    )
    non_operating_expense = sum(
        _normalize_decimal(amounts_data.get(code)) for code in ("NOE_INT", "NOE_OTH")
    )
    net_income = _normalize_decimal(cash_data.get("税引後利益"))
    row_notes = {
        "売上高": "営業活動による売上総額。",
        "売上原価": "仕入・外注・労務など売上連動の変動費。",
        "粗利": "売上高から売上原価を控除した限界利益。",
        "販管費": "人件費や広告費などの固定・半固定費。",
        "営業利益": "本業で創出した利益。",
        "営業外収益": "補助金や受取利息など本業以外の収益。",
        "営業外費用": "支払利息等の本業以外の費用。",
        "経常利益": "営業利益と営業外収支の合計。",
        "税引後利益": "法人税控除後の最終利益。",
    }
    rows: List[List[str]] = [
        ["項目", "金額", "構成比", "注記"],
        ["売上高", format_amount_with_unit(sales, unit), "100%", row_notes.get("売上高", "")],
        ["売上原価", format_amount_with_unit(cogs, unit), _safe_ratio(cogs, sales), row_notes.get("売上原価", "")],
        [
            "粗利",
            format_amount_with_unit(gross, unit),
            format_ratio(_normalize_decimal(metrics_data.get("gross_margin"))),
            row_notes.get("粗利", ""),
        ],
        ["販管費", format_amount_with_unit(opex, unit), _safe_ratio(opex, sales), row_notes.get("販管費", "")],
        [
            "営業利益",
            format_amount_with_unit(op, unit),
            format_ratio(_normalize_decimal(metrics_data.get("op_margin"))),
            row_notes.get("営業利益", ""),
        ],
        ["営業外収益", format_amount_with_unit(non_operating_income, unit), _safe_ratio(non_operating_income, sales), row_notes.get("営業外収益", "")],
        ["営業外費用", format_amount_with_unit(non_operating_expense, unit), _safe_ratio(non_operating_expense, sales), row_notes.get("営業外費用", "")],
        [
            "経常利益",
            format_amount_with_unit(ord_income, unit),
            format_ratio(_normalize_decimal(metrics_data.get("ord_margin"))),
            row_notes.get("経常利益", ""),
        ],
        ["税引後利益", format_amount_with_unit(net_income, unit), _safe_ratio(net_income, sales), row_notes.get("税引後利益", "")],
    ]
    return rows


def _build_initiatives_table(
    amounts_data: Dict[str, Decimal],
    metrics_data: Dict[str, Decimal],
    cash_data: Dict[str, Decimal],
    unit: str,
    fte: Decimal,
) -> List[List[str]]:
    sales = _normalize_decimal(amounts_data.get("REV"))
    breakeven = _normalize_decimal(metrics_data.get("breakeven"))
    gross_margin = _normalize_decimal(metrics_data.get("gross_margin"))
    operating_margin = _normalize_decimal(metrics_data.get("op_margin"))
    operating_cf = _normalize_decimal(cash_data.get("営業キャッシュフロー"))
    investment_cf = _normalize_decimal(cash_data.get("投資キャッシュフロー"))
    net_cf = _normalize_decimal(cash_data.get("キャッシュ増減"))
    breakeven_gap = sales - breakeven if breakeven else Decimal("0")
    revenue_per_fte = sales / fte if fte and fte > 0 and sales else None

    if breakeven > 0:
        if breakeven_gap >= 0:
            revenue_evidence = f"損益分岐点比 {format_amount_with_unit(breakeven_gap, unit)}の余裕"
            revenue_priority = "中" if breakeven_gap >= sales * Decimal("0.1") else "高"
        else:
            revenue_evidence = f"損益分岐点まで {format_amount_with_unit(abs(breakeven_gap), unit)}の不足"
            revenue_priority = "高"
    else:
        revenue_evidence = "損益分岐点の算出には追加データが必要"
        revenue_priority = "高"
    revenue_goal = "重点チャネルでの獲得効率を高め、トップラインを加速させる"
    if revenue_per_fte is not None:
        revenue_goal += f"（FTEあたり {format_amount_with_unit(revenue_per_fte, unit)}）"

    if gross_margin >= Decimal("0.45"):
        margin_priority = "中"
    elif gross_margin >= Decimal("0.35"):
        margin_priority = "中"
    else:
        margin_priority = "高"
    margin_evidence = f"粗利率 {format_ratio(gross_margin)}"
    margin_goal = "原価・値付け戦略の再設計と高付加価値商品の強化で利益率を底上げ"

    net_after_invest = operating_cf + investment_cf
    if net_after_invest < 0 or net_cf < 0:
        cash_priority = "高"
    else:
        cash_priority = "中"
    cash_evidence = (
        f"営業CF {format_amount_with_unit(operating_cf, unit)} / 投資CF {format_amount_with_unit(investment_cf, unit)}"
    )
    cash_goal = "運転資金と調達計画を精緻化し、キャッシュ創出力を強化"

    return [
        ["施策テーマ", "狙い", "優先度", "根拠指標"],
        ["成長アクセラレーション", revenue_goal, revenue_priority, revenue_evidence],
        ["利益率強化プログラム", margin_goal, margin_priority, margin_evidence],
        ["キャッシュマネジメント強化", cash_goal, cash_priority, cash_evidence],
    ]


def _build_financial_simulation_table(
    bundle,
    cash_data: Dict[str, Decimal],
    metrics_data: Dict[str, Decimal],
    unit: str,
) -> List[List[str]]:
    operating_cf = _normalize_decimal(cash_data.get("営業キャッシュフロー"))
    investment_cf = _normalize_decimal(cash_data.get("投資キャッシュフロー"))
    financing_cf = _normalize_decimal(cash_data.get("財務キャッシュフロー"))
    net_cf = _normalize_decimal(cash_data.get("キャッシュ増減"))
    depreciation = _normalize_decimal(cash_data.get("減価償却"))
    investment_metrics = cash_data.get("investment_metrics") or {}
    payback_years_raw = investment_metrics.get("payback_period_years")
    discount_rate_raw = investment_metrics.get("discount_rate")
    npv = _normalize_decimal(investment_metrics.get("npv"))
    payback_value = "—"
    payback_comment = "投資回収期間の算定には追加シナリオが必要です。"
    if payback_years_raw:
        payback_years = float(_normalize_decimal(payback_years_raw))
        payback_value = f"{payback_years:.1f}年"
        payback_comment = (
            "投資回収は標準的なレンジで完了する見込みです。"
            if payback_years <= 5
            else "投資回収が5年以上となるため、キャッシュ創出施策を強化してください。"
        )
    discount_rate_text = "—"
    if discount_rate_raw is not None:
        discount_rate = _normalize_decimal(discount_rate_raw)
        discount_rate_text = format_ratio(discount_rate)
    npv_comment = (
        f"割引率 {discount_rate_text} 基準でNPVプラスを維持。"
        if npv > 0
        else "NPVがマイナスのため、投資前提の再検証が必要です。"
        if npv < 0
        else "NPVは概ねニュートラルです。"
    )
    capex_items = getattr(getattr(bundle, "capex", None), "items", []) if bundle else []
    capex_total = sum(Decimal(getattr(item, "amount", 0)) for item in capex_items)
    capex_names = [getattr(item, "name", "") for item in capex_items if getattr(item, "name", "")]
    if len(capex_names) > 3:
        capex_comment = "／".join(capex_names[:3]) + f" ほか{len(capex_names) - 3}件"
    else:
        capex_comment = "／".join(capex_names)
    capex_comment = capex_comment or "主要設備投資の詳細は計画タブを参照"
    loan_items = getattr(getattr(bundle, "loans", None), "items", []) if bundle else []
    loan_count = len(loan_items)
    operating_comment = f"減価償却 {format_amount_with_unit(depreciation, unit)} を含むキャッシュ創出力"
    investment_comment = (
        f"初年度投資総額 {format_amount_with_unit(capex_total, unit)}｜{capex_comment}"
    )
    financing_comment = f"借入{loan_count}件の償還スケジュールを反映"
    net_comment = (
        "投資後もキャッシュ増減はプラスで推移します。"
        if net_cf >= 0
        else "キャッシュが減少するため、運転資金確保と調達計画の見直しが必要です。"
    )

    return [
        ["シミュレーション項目", "FY計画値", "コメント"],
        ["営業キャッシュフロー", format_amount_with_unit(operating_cf, unit), operating_comment],
        ["投資キャッシュフロー", format_amount_with_unit(investment_cf, unit), investment_comment],
        ["財務キャッシュフロー", format_amount_with_unit(financing_cf, unit), financing_comment],
        ["キャッシュ増減", format_amount_with_unit(net_cf, unit), net_comment],
        ["投資回収期間", payback_value, payback_comment],
        ["NPV（割引率基準）", format_amount_with_unit(npv, unit), npv_comment],
    ]


def _create_summary_chart(amounts_data: Dict[str, Decimal]) -> bytes:
    categories = ["売上高", "粗利", "営業利益", "経常利益"]
    values = [
        float(Decimal(amounts_data.get(code, Decimal("0"))))
        for code in ("REV", "GROSS", "OP", "ORD")
    ]
    max_value = max(values) if values else 0
    fig, ax = plt.subplots(figsize=(5.4, 3.3))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    if max_value > 0:
        ax.axvspan(0, max_value, color=THEME_COLORS["primary_light"], alpha=0.12)
    bars = ax.barh(
        categories,
        values,
        color=[
            THEME_COLORS["chart_blue"],
            THEME_COLORS["chart_green"],
            THEME_COLORS["chart_orange"],
            THEME_COLORS["chart_purple"],
        ],
        edgecolor=THEME_COLORS["primary"],
    )
    ax.set_xlabel("金額")
    ax.set_xlim(left=0, right=max_value * 1.1 if max_value else 1)
    ax.grid(axis="x", linestyle="--", alpha=0.3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color(THEME_COLORS["neutral"])
    ax.spines["bottom"].set_color(THEME_COLORS["neutral"])
    labels = [f"¥{value:,.0f}" for value in values]
    ax.bar_label(
        bars,
        labels=labels,
        label_type="edge",
        padding=6,
        color=THEME_COLORS["accent"],
        fontsize=10,
        fontweight="bold",
    )
    fig.tight_layout()
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png", dpi=220)
    plt.close(fig)
    buffer.seek(0)
    return buffer.getvalue()


def _ensure_pdf_font() -> str:
    """Register a Japanese-capable font for ReportLab if needed."""

    try:
        pdfmetrics.getFont(PDF_FONT_NAME)
    except KeyError:
        pdfmetrics.registerFont(UnicodeCIDFont(PDF_FONT_NAME))
    return PDF_FONT_NAME


def _create_pdf_report(
    title: str,
    subtitle: str,
    sections: Sequence[str],
    tables: Dict[str, List[List[str]]],
    *,
    chart_bytes: bytes | None,
    logo_bytes: bytes | None,
    seal_bytes: bytes | None,
    section_meta: Dict[str, ReportSectionMeta] | None = None,
    section_notes: Dict[str, str] | None = None,
    include_notes: bool = False,
) -> bytes:
    font_name = _ensure_pdf_font()
    buffer = io.BytesIO()
    document = SimpleDocTemplate(buffer, pagesize=A4, title=title)
    styles = getSampleStyleSheet()
    heading_style = ParagraphStyle(
        "HeadingJP",
        parent=styles["Heading2"],
        fontName=font_name,
        textColor=colors.HexColor(THEME_COLORS["primary"]),
    )
    title_style = ParagraphStyle(
        "TitleJP",
        parent=styles["Title"],
        fontName=font_name,
        fontSize=18,
        leading=22,
    )
    subtitle_style = ParagraphStyle(
        "SubtitleJP",
        parent=styles["Normal"],
        fontName=font_name,
        fontSize=11,
        textColor=colors.HexColor(THEME_COLORS["text_subtle"]),
    )
    body_style = ParagraphStyle(
        "BodyJP",
        parent=styles["Normal"],
        fontName=font_name,
        fontSize=11,
        leading=16,
    )
    note_style = ParagraphStyle(
        "NoteJP",
        parent=styles["Normal"],
        fontName=font_name,
        fontSize=9,
        leading=12,
        textColor=colors.HexColor(THEME_COLORS["text_subtle"]),
    )
    meta_lookup = section_meta or {}
    notes_lookup = section_notes or {}
    story: List = []
    if logo_bytes:
        story.append(PlatypusImage(io.BytesIO(logo_bytes), width=140, height=42))
        story.append(Spacer(1, 12))
    story.append(Paragraph(title, title_style))
    story.append(Paragraph(subtitle, subtitle_style))
    story.append(Spacer(1, 18))
    if chart_bytes:
        story.append(PlatypusImage(io.BytesIO(chart_bytes), width=420, height=240))
        story.append(Spacer(1, 18))
    for section in sections:
        table_data = tables.get(section)
        if not table_data:
            continue
        meta = meta_lookup.get(section)
        heading_text = meta.title if meta else section
        story.append(Paragraph(heading_text, heading_style))
        if meta and meta.description:
            story.append(Paragraph(meta.description, note_style))
        story.append(Spacer(1, 6))
        col_count = len(table_data[0]) if table_data else 0
        if col_count == 2:
            col_widths = [220, 180]
        elif col_count == 3:
            col_widths = [160, 120, 120]
        elif col_count == 4:
            col_widths = [150, 120, 90, 130]
        elif col_count == 5:
            col_widths = [140, 110, 80, 80, 130]
        else:
            col_widths = None
        table_component = Table(table_data, colWidths=col_widths)
        table_styles = [
            ("FONTNAME", (0, 0), (-1, -1), font_name),
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(THEME_COLORS["surface_alt"])),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor(THEME_COLORS["primary"])),
            ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
            ("ALIGN", (0, 0), (0, -1), "LEFT"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor(THEME_COLORS["neutral"])),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.HexColor(THEME_COLORS["surface"])],),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]
        if col_count >= 3:
            table_styles.append(("ALIGN", (col_count - 1, 1), (col_count - 1, -1), "LEFT"))
        table_component.setStyle(TableStyle(table_styles))
        story.append(table_component)
        if include_notes:
            note_text = notes_lookup.get(section)
            if note_text:
                story.append(Spacer(1, 6))
                story.append(Paragraph(f"※{note_text}", note_style))
        story.append(Spacer(1, 16))
    if seal_bytes:
        story.append(Spacer(1, 12))
        story.append(Paragraph("承認", body_style))
        story.append(Spacer(1, 6))
        story.append(PlatypusImage(io.BytesIO(seal_bytes), width=120, height=120))
    document.build(story)
    buffer.seek(0)
    return buffer.getvalue()


def _create_excel_report(
    sections: Sequence[str],
    tables: Dict[str, List[List[str]]],
    *,
    chart_bytes: bytes | None,
    logo_bytes: bytes | None,
    metadata: Dict[str, str],
) -> bytes:
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        section_to_sheet = {
            "損益計画": "損益計画",
            "資金繰り表": "資金繰り",
            "投資計画": "投資計画",
        }
        for section in sections:
            table = tables.get(section)
            if not table:
                continue
            df = pd.DataFrame(table[1:], columns=table[0])
            sheet_name = section_to_sheet.get(section, section[:31])
            df.to_excel(writer, sheet_name=sheet_name, index=False)
        meta_df = pd.DataFrame(
            {"項目": list(metadata.keys()), "値": list(metadata.values())}
        )
        meta_df.to_excel(writer, sheet_name="メタ情報", index=False)
        workbook = writer.book
        if logo_bytes and "損益計画" in workbook.sheetnames:
            image = ExcelImage(io.BytesIO(logo_bytes))
            image.width = 200
            image.height = 70
            workbook["損益計画"].add_image(image, "E2")
        if chart_bytes and "損益計画" in workbook.sheetnames:
            chart_image = ExcelImage(io.BytesIO(chart_bytes))
            chart_image.width = 520
            chart_image.height = 320
            workbook["損益計画"].add_image(chart_image, "H2")
    buffer.seek(0)
    return buffer.getvalue()


def _create_word_report(
    title: str,
    subtitle: str,
    sections: Sequence[str],
    tables: Dict[str, List[List[str]]],
    *,
    chart_bytes: bytes | None,
    logo_bytes: bytes | None,
    seal_bytes: bytes | None,
    section_meta: Dict[str, ReportSectionMeta] | None = None,
    section_notes: Dict[str, str] | None = None,
    include_notes: bool = False,
) -> bytes:
    doc = Document()
    meta_lookup = section_meta or {}
    notes_lookup = section_notes or {}
    if logo_bytes:
        doc.add_picture(io.BytesIO(logo_bytes), width=Inches(1.8))
    doc.add_heading(title, level=1)
    doc.add_paragraph(subtitle)
    if chart_bytes:
        doc.add_picture(io.BytesIO(chart_bytes), width=Inches(5.5))
    for section in sections:
        table = tables.get(section)
        if not table:
            continue
        meta = meta_lookup.get(section)
        heading_text = meta.title if meta else section
        doc.add_heading(heading_text, level=2)
        if meta and meta.description:
            doc.add_paragraph(meta.description)
        word_table = doc.add_table(rows=len(table), cols=len(table[0]))
        word_table.style = "Light List Accent 1"
        for row_idx, row in enumerate(table):
            for col_idx, value in enumerate(row):
                word_table.rows[row_idx].cells[col_idx].text = str(value)
        if include_notes:
            note_text = notes_lookup.get(section)
            if note_text:
                doc.add_paragraph(f"※{note_text}")
    if seal_bytes:
        doc.add_paragraph("承認")
        doc.add_picture(io.BytesIO(seal_bytes), width=Inches(1.5))
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def _create_powerpoint_report(
    title: str,
    subtitle: str,
    sections: Sequence[str],
    tables: Dict[str, List[List[str]]],
    *,
    chart_bytes: bytes | None,
    logo_bytes: bytes | None,
    section_meta: Dict[str, ReportSectionMeta] | None = None,
    section_notes: Dict[str, str] | None = None,
) -> bytes:
    prs = Presentation()
    meta_lookup = section_meta or {}
    notes_lookup = section_notes or {}

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = subtitle
    if logo_bytes:
        title_slide.shapes.add_picture(io.BytesIO(logo_bytes), PptxInches(8.3), PptxInches(0.35), width=PptxInches(1.6))
    if chart_bytes:
        title_slide.shapes.add_picture(io.BytesIO(chart_bytes), PptxInches(0.7), PptxInches(2.1), width=PptxInches(8.6))

    for section in sections:
        table_data = tables.get(section)
        meta = meta_lookup.get(section)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = meta.title if meta else section
        left = PptxInches(0.6)
        top = PptxInches(1.5)
        if meta and meta.description:
            desc_box = slide.shapes.add_textbox(left, PptxInches(0.9), PptxInches(9.0), PptxInches(0.8))
            desc_frame = desc_box.text_frame
            desc_frame.text = meta.description
            for paragraph in desc_frame.paragraphs:
                paragraph.font.size = PptxPt(14)
                paragraph.font.color.rgb = RGBColor(75, 93, 117)
                paragraph.space_after = 0
                paragraph.space_before = 0
        table_height = PptxInches(0.5)
        if table_data:
            rows = len(table_data)
            cols = len(table_data[0])
            table_height = PptxInches(0.5 + 0.35 * rows)
            table_shape = slide.shapes.add_table(rows, cols, left, top, PptxInches(9.0), table_height).table
            for i, row in enumerate(table_data):
                for j, value in enumerate(row):
                    cell = table_shape.cell(i, j)
                    cell.text = str(value)
                    cell.text_frame.word_wrap = True
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = PptxPt(13 if i == 0 else 12)
                        paragraph.font.bold = i == 0
                        paragraph.font.color.rgb = RGBColor(12, 29, 51) if i != 0 else RGBColor(255, 255, 255)
                        paragraph.space_after = 0
                        paragraph.space_before = 0
            table_shape.first_row = True
            for header_cell in table_shape.rows[0].cells:
                header_cell.fill.solid()
                header_cell.fill.fore_color.rgb = RGBColor(28, 106, 216)
        if notes_lookup and notes_lookup.get(section):
            note_box = slide.shapes.add_textbox(
                left,
                top + table_height + PptxInches(0.2),
                PptxInches(9.0),
                PptxInches(0.7),
            )
            note_frame = note_box.text_frame
            note_frame.text = f"※{notes_lookup[section]}"
            for paragraph in note_frame.paragraphs:
                paragraph.font.size = PptxPt(12)
                paragraph.font.color.rgb = RGBColor(75, 93, 117)
                paragraph.space_after = 0
                paragraph.space_before = 0

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def _execute_with_spinner(label: str, task: Callable[[], T]) -> T | None:
    """Run a task while showing a spinner and handle failures gracefully."""

    try:
        with st.spinner(f"{label}を生成しています..."):
            return task()
    except Exception:
        st.error(
            f"{label}の生成に失敗しました。入力内容を見直して再度お試しください。"
            f" 解決しない場合は {SUPPORT_CONTACT} までお問い合わせください。"
        )
        return None

default_report_options = {
    "selected_formats": ["PDF", "PowerPoint", "Excel", "Word"],
    "title": f"経営計画サマリー FY{fiscal_year}",
    "subtitle": f"表示単位: {unit} ｜ FTE: {fte}",
    "include_charts": True,
    "sections": list(REPORT_SECTIONS),
    "include_logo": True,
    "include_seal": False,
    "include_notes": True,
}
report_options = st.session_state.setdefault("report_options", default_report_options.copy())
for key, value in default_report_options.items():
    report_options.setdefault(key, value)

with st.form("report_options_form"):
    st.markdown("### レポート設定")
    selected_formats = st.multiselect(
        "出力フォーマット",
        ["PDF", "PowerPoint", "Excel", "Word"],
        default=report_options.get("selected_formats", ["PDF", "PowerPoint", "Excel", "Word"]),
    )
    title_input = st.text_input(
        "レポートタイトル",
        value=report_options.get("title", default_report_options["title"]),
    )
    subtitle_input = st.text_input(
        "サブタイトル",
        value=report_options.get("subtitle", default_report_options["subtitle"]),
    )
    selected_sections = st.multiselect(
        "含めるセクション",
        list(REPORT_SECTIONS),
        default=report_options.get("sections", list(REPORT_SECTIONS)),
        format_func=lambda key: REPORT_SECTION_LOOKUP.get(key, ReportSectionMeta(key, key, "")).title,
    )
    for section_key in selected_sections:
        meta = REPORT_SECTION_LOOKUP.get(section_key)
        if meta and meta.description:
            st.caption(f"・{meta.title} — {meta.description}")
    include_charts = st.checkbox(
        "KPIハイライトの図表を含める",
        value=report_options.get("include_charts", True),
    )
    include_logo = st.checkbox(
        "ヘッダーに企業ロゴを表示",
        value=report_options.get("include_logo", True),
    )
    include_seal = st.checkbox(
        "印影画像を挿入する",
        value=report_options.get("include_seal", False),
    )
    include_notes = st.toggle(
        "金融機関向け注記を併記する",
        value=report_options.get("include_notes", True),
        help="現状分析・施策提案・財務シミュレーションの表に注釈と補足を含めます。",
    )
    submitted = st.form_submit_button("設定を更新", type="primary")
    if submitted:
        report_options.update(
            {
                "selected_formats": selected_formats or ["PDF"],
                "title": title_input.strip() or default_report_options["title"],
                "subtitle": subtitle_input.strip() or default_report_options["subtitle"],
                "sections": selected_sections or list(REPORT_SECTIONS),
                "include_charts": include_charts,
                "include_logo": include_logo,
                "include_seal": include_seal,
                "include_notes": include_notes,
            }
        )
        st.session_state["report_options"] = report_options
        st.success("レポート設定を更新しました。", icon="✔")

logo_upload = st.file_uploader("企業ロゴ (PNG/JPG/SVG)", type=["png", "jpg", "jpeg", "svg"], key="report_logo_upload")
if logo_upload is not None:
    st.session_state["report_logo_bytes"] = logo_upload.getvalue()
    st.session_state["report_logo_name"] = logo_upload.name
    st.toast("ロゴを読み込みました。", icon="図")

seal_upload = st.file_uploader("印影画像 (PNG/JPG)", type=["png", "jpg", "jpeg"], key="report_seal_upload")
if seal_upload is not None:
    st.session_state["report_seal_bytes"] = seal_upload.getvalue()
    st.session_state["report_seal_name"] = seal_upload.name
    st.toast("印影を読み込みました。", icon="印")

logo_bytes = (
    st.session_state.get("report_logo_bytes") if report_options.get("include_logo", True) else None
)
seal_bytes = (
    st.session_state.get("report_seal_bytes") if report_options.get("include_seal", False) else None
)

if logo_bytes:
    st.image(logo_bytes, width=160, caption="ヘッダーロゴ プレビュー")
if seal_bytes:
    st.image(seal_bytes, width=120, caption="印影プレビュー")

selected_sections = [section for section in report_options.get("sections", REPORT_SECTIONS) if section in REPORT_SECTIONS]
if not selected_sections:
    selected_sections = list(REPORT_SECTIONS)

include_notes_flag = bool(report_options.get("include_notes", True))
section_notes = {key: meta.note for key, meta in REPORT_SECTION_LOOKUP.items() if meta.note}

tables = {
    "エグゼクティブサマリー": _build_executive_summary_table(amounts, metrics, cash_flow_data, unit, fte),
    "現状分析": _build_current_analysis_table(amounts, metrics, cash_flow_data, unit),
    "施策提案": _build_initiatives_table(amounts, metrics, cash_flow_data, unit, fte),
    "財務シミュレーション": _build_financial_simulation_table(bundle, cash_flow_data, metrics, unit),
}
chart_bytes = _create_summary_chart(amounts) if report_options.get("include_charts", True) else None
metadata = {
    "作成日時": datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
    "会計年度": f"FY{fiscal_year}",
    "表示単位": unit,
    "FTE": str(fte),
    "注記付与": "あり" if include_notes_flag else "なし",
    "テンプレート": "McKinsey Inspired",
}

if include_notes_flag:
    st.caption("注記付きのMcKinseyスタイルテンプレートで、金融機関提出資料に対応します。")
else:
    st.caption("注記なしのシンプルテンプレートで出力します。")

pdf_tab, excel_tab, word_tab, ppt_tab = st.tabs(["PDF", "Excel", "Word", "PowerPoint"])

with pdf_tab:
    st.subheader("PDFレポート")
    if "PDF" not in report_options.get("selected_formats", []):
        st.info("PDF出力はオフになっています。")
    else:
        pdf_bytes = _execute_with_spinner(
            "PDFレポート",
            lambda: _create_pdf_report(
                report_options.get("title", default_report_options["title"]),
                report_options.get("subtitle", default_report_options["subtitle"]),
                selected_sections,
                tables,
                chart_bytes=chart_bytes,
                logo_bytes=logo_bytes,
                seal_bytes=seal_bytes,
                section_meta=REPORT_SECTION_LOOKUP,
                section_notes=section_notes if include_notes_flag else None,
                include_notes=include_notes_flag,
            ),
        )
        if pdf_bytes is not None:
            st.download_button(
                "［PDF］ダウンロード",
                data=pdf_bytes,
                file_name=f"plan_report_{fiscal_year}.pdf",
                mime="application/pdf",
                **use_container_width_kwargs(st.download_button),
            )

with excel_tab:
    st.subheader("Excelレポート")
    if "Excel" not in report_options.get("selected_formats", []):
        st.info("Excel出力はオフになっています。")
    else:
        excel_bytes = _execute_with_spinner(
            "Excelレポート",
            lambda: _create_excel_report(
                selected_sections,
                tables,
                chart_bytes=chart_bytes,
                logo_bytes=logo_bytes,
                metadata=metadata,
            ),
        )
        if excel_bytes is not None:
            st.download_button(
                "［Excel］ダウンロード",
                data=excel_bytes,
                file_name=f"plan_report_{fiscal_year}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                **use_container_width_kwargs(st.download_button),
            )

with word_tab:
    st.subheader("Wordレポート")
    if "Word" not in report_options.get("selected_formats", []):
        st.info("Word出力はオフになっています。")
    else:
        word_bytes = _execute_with_spinner(
            "Wordレポート",
            lambda: _create_word_report(
                report_options.get("title", default_report_options["title"]),
                report_options.get("subtitle", default_report_options["subtitle"]),
                selected_sections,
                tables,
                chart_bytes=chart_bytes,
                logo_bytes=logo_bytes,
                seal_bytes=seal_bytes,
                section_meta=REPORT_SECTION_LOOKUP,
                section_notes=section_notes if include_notes_flag else None,
                include_notes=include_notes_flag,
            ),
        )
        if word_bytes is not None:
            st.download_button(
                "［Word］ダウンロード",
                data=word_bytes,
                file_name=f"plan_report_{fiscal_year}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                **use_container_width_kwargs(st.download_button),
            )

with ppt_tab:
    st.subheader("PowerPointレポート")
    if "PowerPoint" not in report_options.get("selected_formats", []):
        st.info("PowerPoint出力はオフになっています。")
    else:
        ppt_bytes = _execute_with_spinner(
            "PowerPointレポート",
            lambda: _create_powerpoint_report(
                report_options.get("title", default_report_options["title"]),
                report_options.get("subtitle", default_report_options["subtitle"]),
                selected_sections,
                tables,
                chart_bytes=chart_bytes,
                logo_bytes=logo_bytes,
                section_meta=REPORT_SECTION_LOOKUP,
                section_notes=section_notes if include_notes_flag else None,
            ),
        )
        if ppt_bytes is not None:
            st.download_button(
                "［PowerPoint］ダウンロード",
                data=ppt_bytes,
                file_name=f"plan_report_{fiscal_year}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                **use_container_width_kwargs(st.download_button),
            )
